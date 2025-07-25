import streamlit as st
import pandas as pd
import numpy as np
import json
from datetime import datetime
import plotly.express as px
import plotly.graph_objects as go
import io

# Optional imports for exports - gracefully handle if not available
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from reportlab.lib.pagesizes import A4, letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.graphics.shapes import Drawing
    from reportlab.graphics.charts.barcharts import VerticalBarChart
    from reportlab.graphics.charts.piecharts import Pie
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

# Configure Streamlit page
st.set_page_config(
    page_title="HR ROI Calculator",
    page_icon="🎯",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for better styling
st.markdown("""
<style>
    .metric-card {
        background: white;
        padding: 1rem;
        border-radius: 0.5rem;
        box-shadow: 0 1px 3px rgba(0,0,0,0.1);
        text-align: center;
    }
    .initiative-card {
        background: #f8f9fa;
        padding: 1rem;
        border-radius: 0.5rem;
        border-left: 4px solid #007bff;
        margin: 1rem 0;
    }
</style>
""", unsafe_allow_html=True)

# Initiative Templates - Focus on individual HR programs
INITIATIVE_TEMPLATES = {
    'leadership_development': {
        'name': "Leadership Development Program",
        'description': "Comprehensive leadership training with coaching and assessments - improves productivity, retention, team performance, and reduces sick leave",
        'participants': 25,
        'program_duration': 6,
        'avg_salary': 100000,
        'facilitator_costs': 80000,
        'materials_costs': 20000,
        'venue_costs': 30000,
        'travel_costs': 20000,
        'productivity_gain': 18,
        'retention_improvement': 30,
        'team_performance_gain': 15,
        'sick_leave_reduction': 20,
        'typical_roi': "500-1000%"
    },
    'executive_coaching': {
        'name': "Executive Coaching Initiative", 
        'description': "1-on-1 coaching for senior leaders - reduces stress, improves decision-making, and creates healthier work environments",
        'participants': 10,
        'program_duration': 12,
        'avg_salary': 150000,
        'facilitator_costs': 120000,
        'materials_costs': 5000,
        'venue_costs': 0,
        'travel_costs': 5000,
        'productivity_gain': 25,
        'retention_improvement': 40,
        'team_performance_gain': 20,
        'sick_leave_reduction': 25,
        'typical_roi': "600-1200%"
    },
    'time_to_fill_optimization': {
        'name': "Time to Fill Optimization",
        'description': "Reduce time to fill vacant positions through process improvements, technology, and recruiter training - includes both cost savings and revenue impact analysis",
        'annual_positions': 30,
        'current_time_to_fill': 60,
        'target_time_to_fill': 35,
        'avg_position_salary': 85000,
        'productivity_loss_rate': 70,
        'overtime_multiplier': 1.5,
        'team_impact_factor': 15,
        'optimization_investment': 45000,
        'training_costs': 15000,
        'technology_costs': 25000,
        'revenue_generating_percentage': 60,
        'revenue_per_employee_daily': 800,
        'customer_impact_factor': 25,
        'typical_roi': "400-800%"
    },
    'onboarding_excellence': {
        'name': "Structured Onboarding Program",
        'description': "Comprehensive new hire integration and training",
        'annual_new_hires': 50,
        'current_time_to_productivity': 4,
        'productivity_acceleration': 50,
        'new_hire_retention_rate': 75,
        'onboarding_retention_improvement': 20,
        'onboarding_program_cost': 40000,
        'typical_roi': "300-450%"
    },
    'engagement_retention': {
        'name': "Employee Engagement & Retention",
        'description': "Surveys, feedback systems, and engagement initiatives",
        'total_employees': 500,
        'current_engagement_score': 6.5,
        'current_turnover': 18,
        'engagement_improvement': 2.0,
        'retention_improvement': 30,
        'engagement_program_cost': 60000,
        'typical_roi': "250-400%"
    },
    'talent_development': {
        'name': "Internal Talent Development",
        'description': "Skills training, mentoring, and succession planning",
        'development_participants': 100,
        'internal_mobility_rate': 60,
        'succession_readiness': 40,
        'mobility_improvement': 25,
        'succession_improvement': 35,
        'development_program_cost': 100000,
        'typical_roi': "200-300%"
    }
}

def format_currency(amount):
    """Format amount as currency"""
    return f"${amount:,.0f}"

def get_roi_status(roi):
    """Get status and color for ROI"""
    if roi >= 500:
        return "🟢 Exceptional (500%+)"
    elif roi >= 300:
        return "🟢 Excellent (300-499%)"
    elif roi >= 200:
        return "🟡 Good (200-299%)"
    elif roi >= 100:
        return "🟠 Moderate (100-199%)"
    else:
        return "🔴 Needs Review (<100%)"

def calculate_leadership_roi(params):
    """Calculate Leadership Development ROI - FIXED VERSION using only incremental costs"""
    
    # ONLY TRUE INCREMENTAL COSTS (no salary/opportunity costs)
    total_incremental_costs = (
        params['facilitator_costs'] + 
        params['materials_costs'] + 
        params['venue_costs'] + 
        params.get('travel_costs', 20000)
    )
    
    # INCREMENTAL BENEFITS (annual value above baseline)
    productivity_benefit = params['participants'] * params['avg_salary'] * (params['productivity_gain'] / 100)
    
    current_turnover = params.get('current_turnover', 18)
    replacement_cost = params.get('replacement_cost', 1.5)
    retention_savings = (
        params['participants'] * (current_turnover / 100) * 
        (params['retention_improvement'] / 100) * 
        params['avg_salary'] * replacement_cost
    )
    
    team_size = params.get('team_size', 8)
    team_benefit = (
        params['participants'] * team_size * 
        (params['avg_salary'] * 0.7) * (params['team_performance_gain'] / 100)
    )
    
    # NEW: Sick leave reduction benefit
    current_sick_days = params.get('current_sick_days', 8)
    sick_leave_reduction_pct = params.get('sick_leave_reduction', 20) / 100
    daily_salary = params['avg_salary'] / 250
    
    # Both direct (participants) and indirect (team members) benefit from better leadership
    total_affected_employees = params['participants'] + (params['participants'] * team_size)
    sick_leave_savings = (
        total_affected_employees * 
        current_sick_days * 
        sick_leave_reduction_pct * 
        daily_salary * 1.3
    )
    
    total_annual_benefits = productivity_benefit + retention_savings + team_benefit + sick_leave_savings
    
    # Calculate ROI
    roi = ((total_annual_benefits - total_incremental_costs) / total_incremental_costs * 100) if total_incremental_costs > 0 else 0
    payback_months = (total_incremental_costs / (total_annual_benefits / 12)) if total_annual_benefits > 0 else 0
    
    return {
        'total_costs': total_incremental_costs,
        'annual_benefits': total_annual_benefits,
        'roi': roi,
        'payback_months': payback_months,
        'net_annual_benefit': total_annual_benefits - total_incremental_costs,
        'benefit_breakdown': {
            'productivity': productivity_benefit,
            'retention': retention_savings,
            'team_performance': team_benefit,
            'sick_leave_reduction': sick_leave_savings
        },
        'cost_breakdown': {
            'facilitator_costs': params['facilitator_costs'],
            'materials_costs': params['materials_costs'],
            'venue_costs': params['venue_costs'],
            'travel_costs': params.get('travel_costs', 20000)
        }
    }

def calculate_time_to_fill_roi(params):
    """Calculate Time to Fill Optimization ROI with Revenue Impact Analysis"""
    annual_positions = params['annual_positions']
    current_days = params['current_time_to_fill']
    target_days = params['target_time_to_fill']
    days_saved = current_days - target_days
    
    avg_salary = params['avg_position_salary']
    daily_salary = avg_salary / 250
    
    # === COST SAVINGS ANALYSIS (Original) ===
    # 1. Productivity recovery
    productivity_loss_rate = params['productivity_loss_rate'] / 100
    productivity_recovery = annual_positions * days_saved * daily_salary * productivity_loss_rate
    
    # 2. Overtime cost reduction
    overtime_multiplier = params.get('overtime_multiplier', 1.5)
    overtime_hours_per_day = params.get('overtime_hours_per_day', 2)
    overtime_savings = (
        annual_positions * days_saved * overtime_hours_per_day * 
        (daily_salary / 8) * (overtime_multiplier - 1)
    )
    
    # 3. Team productivity impact
    team_impact_factor = params['team_impact_factor'] / 100
    team_size = params.get('team_size', 6)
    team_productivity_gain = (
        annual_positions * days_saved * team_size * 
        daily_salary * 0.8 * team_impact_factor
    )
    
    # 4. Faster time to productivity
    faster_onboarding_value = annual_positions * (days_saved * 0.3) * daily_salary * 0.6
    
    total_cost_savings = productivity_recovery + overtime_savings + team_productivity_gain + faster_onboarding_value
    
    # === REVENUE IMPACT ANALYSIS (New) ===
    revenue_generating_pct = params.get('revenue_generating_percentage', 60) / 100
    revenue_per_employee_daily = params.get('revenue_per_employee_daily', 800)
    customer_impact_factor = params.get('customer_impact_factor', 25) / 100
    
    # Revenue-generating positions
    revenue_positions = annual_positions * revenue_generating_pct
    
    # 1. Direct revenue loss from vacant revenue-generating positions
    direct_revenue_loss_prevented = revenue_positions * days_saved * revenue_per_employee_daily
    
    # 2. Customer service/satisfaction impact on remaining revenue
    customer_impact_loss_prevented = (
        annual_positions * days_saved * revenue_per_employee_daily * 
        customer_impact_factor * params.get('customer_base_factor', 0.15)
    )
    
    # 3. Lost opportunity costs (deals not closed, projects delayed)
    opportunity_cost_factor = params.get('opportunity_cost_factor', 0.20)
    opportunity_loss_prevented = (
        revenue_positions * days_saved * revenue_per_employee_daily * opportunity_cost_factor
    )
    
    # 4. Market share protection (competitors don't gain ground)
    market_share_protection = (
        revenue_positions * days_saved * revenue_per_employee_daily * 
        params.get('market_share_factor', 0.05)
    )
    
    total_revenue_protection = (
        direct_revenue_loss_prevented + customer_impact_loss_prevented + 
        opportunity_loss_prevented + market_share_protection
    )
    
    # === COMBINED BENEFITS ===
    total_annual_benefits = total_cost_savings + total_revenue_protection
    
    # Investment costs
    total_investment = (
        params['optimization_investment'] + 
        params.get('training_costs', 15000) + 
        params.get('technology_costs', 25000)
    )
    
    # Calculate ROI
    roi = ((total_annual_benefits - total_investment) / total_investment * 100) if total_investment > 0 else 0
    payback_months = (total_investment / (total_annual_benefits / 12)) if total_annual_benefits > 0 else 0
    
    # Calculate daily impact metrics
    daily_revenue_at_risk = annual_positions * revenue_generating_pct * revenue_per_employee_daily
    total_revenue_at_risk_current = daily_revenue_at_risk * current_days
    total_revenue_at_risk_target = daily_revenue_at_risk * target_days
    annual_revenue_at_risk_reduction = daily_revenue_at_risk * days_saved * annual_positions
    
    return {
        'total_investment': total_investment,
        'annual_benefits': total_annual_benefits,
        'roi': roi,
        'payback_months': payback_months,
        'days_saved_per_position': days_saved,
        'total_days_saved_annually': annual_positions * days_saved,
        
        # Cost savings breakdown
        'cost_savings_breakdown': {
            'productivity_recovery': productivity_recovery,
            'overtime_savings': overtime_savings,
            'team_productivity_gain': team_productivity_gain,
            'faster_onboarding': faster_onboarding_value
        },
        'total_cost_savings': total_cost_savings,
        
        # Revenue impact breakdown
        'revenue_impact_breakdown': {
            'direct_revenue_protection': direct_revenue_loss_prevented,
            'customer_impact_protection': customer_impact_loss_prevented,
            'opportunity_cost_protection': opportunity_loss_prevented,
            'market_share_protection': market_share_protection
        },
        'total_revenue_protection': total_revenue_protection,
        
        # Revenue metrics
        'daily_revenue_at_risk': daily_revenue_at_risk,
        'revenue_positions': revenue_positions,
        'total_revenue_at_risk_current': total_revenue_at_risk_current,
        'total_revenue_at_risk_target': total_revenue_at_risk_target,
        'annual_revenue_at_risk_reduction': annual_revenue_at_risk_reduction,
        
        'investment_breakdown': {
            'optimization_investment': params['optimization_investment'],
            'training_costs': params.get('training_costs', 15000),
            'technology_costs': params.get('technology_costs', 25000)
        }
    }

def calculate_onboarding_roi(params):
    """Calculate Onboarding ROI"""
    annual_hires = params['annual_new_hires']
    current_productivity_time = params['current_time_to_productivity']
    acceleration = params['productivity_acceleration'] / 100
    
    # Productivity improvement
    months_saved = current_productivity_time * acceleration
    productivity_value = annual_hires * months_saved * (params.get('avg_salary', 95000) / 12) * 0.6
    
    # Retention improvement
    retention_improvement = params['onboarding_retention_improvement'] / 100
    retention_value = annual_hires * retention_improvement * params.get('avg_salary', 95000) * 1.5
    
    total_annual_benefits = productivity_value + retention_value
    total_investment = params['onboarding_program_cost'] + params.get('training_costs', 15000)
    
    roi = ((total_annual_benefits - total_investment) / total_investment * 100) if total_investment > 0 else 0
    
    return {
        'total_investment': total_investment,
        'annual_benefits': total_annual_benefits,
        'roi': roi,
        'benefit_breakdown': {
            'productivity_value': productivity_value,
            'retention_value': retention_value
        }
    }

def calculate_engagement_roi(params):
    """Calculate Engagement & Retention ROI"""
    total_employees = params['total_employees']
    current_turnover = params['current_turnover'] / 100
    engagement_improvement = params['engagement_improvement']
    retention_improvement = params['retention_improvement'] / 100
    
    # Benefits
    productivity_boost = total_employees * params.get('avg_salary', 95000) * (engagement_improvement * 0.02)
    turnover_reduction = total_employees * current_turnover * retention_improvement
    turnover_savings = turnover_reduction * params.get('avg_salary', 95000) * 1.5
    
    total_annual_benefits = productivity_boost + turnover_savings
    total_investment = params['engagement_program_cost'] + params.get('survey_costs', 15000)
    
    roi = ((total_annual_benefits - total_investment) / total_investment * 100) if total_investment > 0 else 0
    
    return {
        'total_investment': total_investment,
        'annual_benefits': total_annual_benefits,
        'roi': roi,
        'benefit_breakdown': {
            'productivity_boost': productivity_boost,
            'turnover_savings': turnover_savings
        }
    }

def calculate_development_roi(params):
    """Calculate Talent Development ROI"""
    participants = params['development_participants']
    mobility_improvement = params['mobility_improvement'] / 100
    
    # Benefits
    performance_gains = participants * params.get('avg_salary', 95000) * 0.12
    internal_hiring_savings = participants * mobility_improvement * 25000
    retention_boost = participants * 0.2 * params.get('avg_salary', 95000) * 1.5
    
    total_annual_benefits = performance_gains + internal_hiring_savings + retention_boost
    total_investment = params['development_program_cost']
    
    roi = ((total_annual_benefits - total_investment) / total_investment * 100) if total_investment > 0 else 0
    
    return {
        'total_investment': total_investment,
        'annual_benefits': total_annual_benefits,
        'roi': roi,
        'benefit_breakdown': {
            'performance_gains': performance_gains,
            'hiring_savings': internal_hiring_savings,
            'retention_boost': retention_boost
        }
    }

def create_pdf_report(initiative_results, overall_roi, total_investment, total_benefits, params_data):
    """Create a comprehensive PDF report"""
    if not REPORTLAB_AVAILABLE:
        return None
    
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, topMargin=0.5*inch)
    styles = getSampleStyleSheet()
    story = []
    
    # Check if single initiative
    is_single_initiative = len(initiative_results) == 1
    
    # Title
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=20,
        spaceAfter=30,
        textColor=colors.darkblue
    )
    
    if is_single_initiative:
        story.append(Paragraph(f"HR ROI Analysis - {initiative_results[0]['Initiative']}", title_style))
    else:
        story.append(Paragraph("HR ROI Calculator - Portfolio Report", title_style))
    
    story.append(Paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}", styles['Normal']))
    story.append(Paragraph("Note: ROI calculated using incremental costs vs. incremental benefits", styles['Italic']))
    story.append(Spacer(1, 20))
    
    # Executive Summary
    if is_single_initiative:
        story.append(Paragraph("Executive Summary", styles['Heading2']))
        
        summary_data = [
            ['Metric', 'Value', 'Status'],
            ['Initiative ROI', f"{overall_roi:.0f}%", get_roi_status(overall_roi)],
            ['Total Investment', format_currency(total_investment), ''],
            ['Total Annual Benefits', format_currency(total_benefits), ''],
            ['Net Annual Benefit', format_currency(total_benefits - total_investment), '']
        ]
    else:
        story.append(Paragraph("Executive Summary", styles['Heading2']))
        
        summary_data = [
            ['Metric', 'Value', 'Status'],
            ['Portfolio ROI', f"{overall_roi:.0f}%", get_roi_status(overall_roi)],
            ['Total Investment', format_currency(total_investment), ''],
            ['Total Annual Benefits', format_currency(total_benefits), ''],
            ['Net Annual Benefit', format_currency(total_benefits - total_investment), '']
        ]
    
    summary_table = Table(summary_data, colWidths=[2*inch, 1.5*inch, 2*inch])
    summary_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black)
    ]))
    
    story.append(summary_table)
    story.append(Spacer(1, 20))
    
    # Initiative Details
    story.append(Paragraph("Initiative Breakdown", styles['Heading2']))
    
    for i, initiative in enumerate(initiative_results):
        story.append(Paragraph(f"{i+1}. {initiative['Initiative']}", styles['Heading3']))
        
        init_data = [
            ['Investment', format_currency(initiative['Investment'])],
            ['Annual Benefits', format_currency(initiative['Annual Benefits'])],
            ['ROI', f"{initiative['ROI (%)']:.0f}%"],
            ['Status', get_roi_status(initiative['ROI (%)'])]
        ]
        
        init_table = Table(init_data, colWidths=[2*inch, 2*inch])
        init_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (0, -1), colors.lightgrey),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        
        story.append(init_table)
        story.append(Spacer(1, 12))
    
    # Methodology Note
    story.append(Paragraph("Methodology Note", styles['Heading2']))
    methodology_text = """This analysis uses incremental cost accounting, where only true additional expenses are counted as costs (e.g., facilitator fees, materials, venues). Employee salary costs are not included since employees are paid regardless of training participation. This approach provides a more accurate view of the actual investment required and expected returns."""
    story.append(Paragraph(methodology_text, styles['Normal']))
    story.append(Spacer(1, 20))
    
    # Recommendations
    story.append(Paragraph("Strategic Recommendations", styles['Heading2']))
    
    if is_single_initiative:
        if overall_roi >= 500:
            recommendation = "✅ Exceptional ROI performance. Proceed with immediate implementation. Consider scaling this program to additional employee populations and similar roles."
        elif overall_roi >= 300:
            recommendation = "✅ Excellent ROI performance. Proceed with implementation. Monitor key metrics closely and prepare for potential expansion."
        elif overall_roi >= 200:
            recommendation = "✅ Strong ROI performance. Proceed with implementation. Consider optimizing program design for even better results."
        elif overall_roi >= 100:
            recommendation = "⚠️ Moderate ROI performance. Proceed with implementation but consider optimizing program design or targeting higher-impact participants."
        else:
            recommendation = "❌ Initiative requires optimization. Review assumptions, implementation strategy, and target population before proceeding."
    else:
        if overall_roi >= 500:
            recommendation = "✅ Exceptional portfolio performance. Proceed with full implementation across all initiatives. Consider scaling successful programs and expanding to additional employee populations."
        elif overall_roi >= 300:
            recommendation = "✅ Excellent portfolio performance. Proceed with implementation, prioritizing highest ROI initiatives first. Monitor key metrics closely during rollout."
        elif overall_roi >= 200:
            recommendation = "✅ Strong portfolio performance. Proceed with implementation, prioritizing highest ROI initiatives first. Monitor key metrics closely during rollout."
        elif overall_roi >= 100:
            recommendation = "⚠️ Moderate portfolio performance. Focus on highest ROI initiatives for immediate implementation. Review and optimize lower-performing programs before proceeding."
        else:
            recommendation = "❌ Portfolio requires significant optimization. Focus resources on highest ROI initiatives only. Reassess assumptions and implementation strategies for underperforming programs."
    
    story.append(Paragraph(recommendation, styles['Normal']))
    story.append(Spacer(1, 20))
    
    # Implementation Priority (only for multiple initiatives)
    if not is_single_initiative:
        story.append(Paragraph("Implementation Priority Matrix", styles['Heading3']))
        sorted_initiatives = sorted(initiative_results, key=lambda x: x['ROI (%)'], reverse=True)
        
        priority_data = [['Priority', 'Initiative', 'ROI', 'Recommendation']]
        for i, init in enumerate(sorted_initiatives):
            if init['ROI (%)'] >= 400:
                priority = "Phase 1 (Immediate)"
            elif init['ROI (%)'] >= 200:
                priority = "Phase 2 (3-6 months)"
            else:
                priority = "Phase 3 (Review)"
            
            priority_data.append([
                priority,
                init['Initiative'],
                f"{init['ROI (%)']:.0f}%",
                "Implement" if init['ROI (%)'] >= 200 else "Optimize"
            ])
        
        priority_table = Table(priority_data, colWidths=[1.5*inch, 2.5*inch, 1*inch, 1*inch])
        priority_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.darkblue),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 10),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.lightblue),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        
        story.append(priority_table)
    else:
        # For single initiatives, add specific implementation guidance
        story.append(Paragraph("Implementation Guidance", styles['Heading3']))
        
        if overall_roi >= 300:
            implementation_guidance = """Recommended Implementation Steps:

1. Secure executive sponsorship and budget approval
2. Develop detailed project plan with clear timelines
3. Establish baseline metrics for measurement
4. Launch pilot program with key stakeholders
5. Monitor progress and adjust as needed
6. Plan for scaling based on results"""
        else:
            implementation_guidance = """Before Implementation:

1. Review and validate assumptions with stakeholders
2. Consider optimizing program design or targeting
3. Establish clear success criteria and measurement plan
4. Start with a smaller pilot to test effectiveness
5. Gather additional data to strengthen business case"""
        
        story.append(Paragraph(implementation_guidance, styles['Normal']))
        story.append(Spacer(1, 20))
    
    doc.build(story)
    buffer.seek(0)
    return buffer

def create_powerpoint_presentation(initiative_results, overall_roi, total_investment, total_benefits):
    """Create a PowerPoint presentation"""
    if not PPTX_AVAILABLE:
        return None
    
    prs = Presentation()
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "HR ROI Calculator Results"
    subtitle.text = f"Portfolio Analysis (Incremental Cost Method)\nGenerated: {datetime.now().strftime('%B %d, %Y')}"
    
    # Slide 2: Executive Summary
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Executive Summary"
    
    summary_text = f"""Portfolio Performance Overview:

• Total ROI: {overall_roi:.0f}%
• Total Investment: {format_currency(total_investment)}
• Total Annual Benefits: {format_currency(total_benefits)}
• Net Annual Benefit: {format_currency(total_benefits - total_investment)}

Status: {get_roi_status(overall_roi)}

Note: Analysis uses incremental cost accounting methodology
"""
    
    content.text = summary_text
    
    # Slide 3: Initiative Comparison
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Initiative Performance Comparison"
    
    # Sort initiatives by ROI
    sorted_initiatives = sorted(initiative_results, key=lambda x: x['ROI (%)'], reverse=True)
    
    comparison_text = "Initiative Rankings:\n\n"
    for i, init in enumerate(sorted_initiatives):
        status_emoji = "🟢" if init['ROI (%)'] >= 500 else "🟢" if init['ROI (%)'] >= 300 else "🟡" if init['ROI (%)'] >= 200 else "🔴"
        comparison_text += f"{i+1}. {init['Initiative']}\n"
        comparison_text += f"   ROI: {init['ROI (%)']:.0f}% {status_emoji}\n"
        comparison_text += f"   Investment: {format_currency(init['Investment'])}\n\n"
    
    content.text = comparison_text
    
    # Slide 4: Methodology
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Methodology: Incremental Cost Accounting"
    
    methodology_text = """Key Principles:

• Costs: Only true incremental expenses counted
  - Facilitator fees, materials, venues, travel
  - Employee salaries excluded (paid regardless)

• Benefits: Incremental value above baseline
  - Productivity improvements
  - Retention savings
  - Team performance gains

• Result: More accurate ROI reflecting actual investment
  vs. expected returns"""
    
    content.text = methodology_text
    
    # Slide 5: Implementation Roadmap
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Implementation Roadmap"
    
    roadmap_text = """Phase 1 (Immediate - 0-3 months):
• Launch initiatives with ROI ≥ 400%
• Secure executive sponsorship
• Establish measurement frameworks

Phase 2 (Short-term - 3-6 months):
• Implement initiatives with ROI 200-399%
• Monitor Phase 1 results
• Adjust programs based on early feedback

Phase 3 (Long-term - 6+ months):
• Review and optimize underperforming initiatives
• Scale successful programs
• Develop next generation of HR initiatives"""
    
    content.text = roadmap_text
    
    # Slide 6: Key Success Factors
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Key Success Factors"
    
    success_text = """Critical Elements for Success:

• Executive Leadership Support
  - Visible commitment from senior leadership
  - Adequate resource allocation

• Robust Measurement & Analytics
  - Clear KPIs and success metrics
  - Regular progress monitoring

• Change Management
  - Comprehensive communication strategy
  - Employee engagement and buy-in

• Continuous Improvement
  - Regular program evaluation
  - Agile adjustment of strategies"""
    
    content.text = success_text
    
    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

def main():
    # Header
    st.markdown("""
    <div style='background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); 
                padding: 2rem; border-radius: 10px; margin-bottom: 2rem;'>
        <h1 style='color: white; margin: 0; font-size: 2.5rem;'>🎯 HR ROI Calculator</h1>
        <p style='color: rgba(255,255,255,0.8); margin: 0.5rem 0 0 0; font-size: 1.2rem;'>
            Calculate ROI for Individual HR Initiatives (Using Incremental Cost Method)
        </p>
    </div>
    """, unsafe_allow_html=True)
    
    # KPI Helper Calculator
    with st.expander("🧮 KPI Helper Calculator - Calculate Your Baseline Metrics", expanded=False):
        st.markdown("**Use these calculators to determine the input values for your ROI analysis:**")
        
        tab1, tab2, tab3, tab4 = st.tabs(["💰 Cost Metrics", "⏰ Time Metrics", "📊 Performance Metrics", "💵 Revenue Metrics"])
        
        with tab1:
            st.subheader("💰 Cost Calculation Helpers")
            
            col1, col2 = st.columns(2)
            with col1:
                st.markdown("**Cost per Hire Calculator**")
                external_costs = st.number_input("External Costs (agencies, job boards) ($)", min_value=0, value=3000, step=100, key="cph_external")
                internal_costs = st.number_input("Internal Costs (recruiter time, interviews) ($)", min_value=0, value=2000, step=100, key="cph_internal")
                total_hires = st.number_input("Number of Hires", min_value=1, value=10, key="cph_hires")
                cost_per_hire = (external_costs + internal_costs) / total_hires
                st.success(f"**Cost per Hire: {format_currency(cost_per_hire)}**")
                
                st.markdown("**Turnover Cost Calculator**")
                avg_salary = st.number_input("Average Salary ($)", min_value=0, value=75000, step=5000, key="tc_salary")
                turnover_multiplier = st.selectbox("Turnover Cost Multiplier", [0.5, 0.75, 1.0, 1.5, 2.0], index=3, key="tc_mult")
                st.info(f"0.5x = Entry level, 1.5x = Professional, 2.0x = Senior roles")
                turnover_cost = avg_salary * turnover_multiplier
                st.success(f"**Turnover Cost: {format_currency(turnover_cost)}**")
            
            with col2:
                st.markdown("**Training Cost Calculator**")
                facilitator_daily_rate = st.number_input("Facilitator Daily Rate ($)", min_value=0, value=2000, step=100, key="tc_facilitator")
                training_days = st.number_input("Training Days", min_value=1, value=5, key="tc_days")
                materials_per_person = st.number_input("Materials per Person ($)", min_value=0, value=500, step=50, key="tc_materials")
                venue_daily = st.number_input("Venue Cost per Day ($)", min_value=0, value=800, step=100, key="tc_venue")
                participants = st.number_input("Number of Participants", min_value=1, value=20, key="tc_participants")
                
                total_training_cost = (facilitator_daily_rate * training_days) + (materials_per_person * participants) + (venue_daily * training_days)
                cost_per_participant = total_training_cost / participants
                
                st.success(f"**Total Training Cost: {format_currency(total_training_cost)}**")
                st.success(f"**Cost per Participant: {format_currency(cost_per_participant)}**")
        
        with tab2:
            st.subheader("⏰ Time Calculation Helpers")
            
            col1, col2 = st.columns(2)
            with col1:
                st.markdown("**Time to Productivity Calculator**")
                role_complexity = st.selectbox("Role Complexity", ["Entry Level", "Mid Level", "Senior Level", "Executive"], key="ttp_complexity")
                complexity_months = {"Entry Level": 2, "Mid Level": 4, "Senior Level": 6, "Executive": 9}
                industry_factor = st.selectbox("Industry Learning Curve", ["Low (Tech/Service)", "Medium (Manufacturing)", "High (Healthcare/Finance)"], key="ttp_industry")
                industry_multiplier = {"Low (Tech/Service)": 0.8, "Medium (Manufacturing)": 1.0, "High (Healthcare/Finance)": 1.3}
                
                base_months = complexity_months[role_complexity]
                adjusted_months = base_months * industry_multiplier[industry_factor]
                st.success(f"**Estimated Time to Productivity: {adjusted_months:.1f} months**")
                
                st.markdown("**Daily Productivity Value Calculator**")
                annual_salary = st.number_input("Annual Salary ($)", min_value=0, value=80000, step=5000, key="dpv_salary")
                productivity_multiplier = st.number_input("Productivity Multiplier", min_value=1.0, max_value=5.0, value=2.5, step=0.1, key="dpv_mult")
                st.info("Productivity multiplier: How much value they create vs. their salary")
                
                daily_value = (annual_salary * productivity_multiplier) / 250
                st.success(f"**Daily Productivity Value: {format_currency(daily_value)}**")
            
            with col2:
                st.markdown("**Time to Fill Benchmarks**")
                position_level = st.selectbox("Position Level", ["Entry Level", "Professional", "Manager", "Director+"], key="ttf_level")
                industry_type = st.selectbox("Industry Type", ["Technology", "Healthcare", "Finance", "Manufacturing", "Retail"], key="ttf_industry")
                
                # Benchmark data (industry averages)
                benchmarks = {
                    "Entry Level": {"Technology": 25, "Healthcare": 30, "Finance": 35, "Manufacturing": 28, "Retail": 20},
                    "Professional": {"Technology": 35, "Healthcare": 45, "Finance": 50, "Manufacturing": 40, "Retail": 30},
                    "Manager": {"Technology": 50, "Healthcare": 60, "Finance": 65, "Manufacturing": 55, "Retail": 45},
                    "Director+": {"Technology": 75, "Healthcare": 85, "Finance": 90, "Manufacturing": 80, "Retail": 65}
                }
                
                benchmark_days = benchmarks[position_level][industry_type]
                st.success(f"**Industry Benchmark: {benchmark_days} days**")
                st.info(f"🎯 Good target: {int(benchmark_days * 0.8)} days (-20%)")
                st.info(f"🚀 Excellent target: {int(benchmark_days * 0.6)} days (-40%)")
        
        with tab3:
            st.subheader("📊 Performance Impact Calculators")
            
            col1, col2 = st.columns(2)
            with col1:
                st.markdown("**Engagement Impact Calculator**")
                current_engagement = st.number_input("Current Engagement Score (1-10)", min_value=1.0, max_value=10.0, value=6.5, step=0.1, key="ei_current")
                target_engagement = st.number_input("Target Engagement Score (1-10)", min_value=1.0, max_value=10.0, value=8.0, step=0.1, key="ei_target")
                
                engagement_improvement = target_engagement - current_engagement
                productivity_impact = engagement_improvement * 0.02  # Research shows 2% productivity gain per engagement point
                
                st.success(f"**Engagement Improvement: +{engagement_improvement:.1f} points**")
                st.success(f"**Estimated Productivity Impact: +{productivity_impact:.1%}**")
                
                st.markdown("**Absenteeism Impact Calculator**")
                current_sick_days = st.number_input("Current Sick Days/Employee/Year", min_value=0, value=8, key="ai_current")
                industry_benchmark = st.selectbox("Industry Benchmark", ["Low (5 days)", "Average (8 days)", "High (12 days)"], key="ai_benchmark")
                benchmark_days = {"Low (5 days)": 5, "Average (8 days)": 8, "High (12 days)": 12}[industry_benchmark]
                
                potential_reduction = max(0, current_sick_days - benchmark_days)
                st.success(f"**Potential Reduction: {potential_reduction} days/year**")
                
            with col2:
                st.markdown("**Retention Impact Calculator**")
                current_turnover = st.number_input("Current Turnover Rate (%)", min_value=0.0, max_value=100.0, value=18.0, step=1.0, key="ri_current")
                industry_avg = st.selectbox("Industry Average", ["Low (8%)", "Medium (15%)", "High (25%)"], key="ri_industry")
                industry_turnover = {"Low (8%)": 8, "Medium (15%)": 15, "High (25%)": 25}[industry_avg]
                
                improvement_potential = max(0, current_turnover - industry_turnover)
                st.success(f"**Improvement Potential: -{improvement_potential:.1f}%**")
                
                total_employees = st.number_input("Total Employees", min_value=1, value=100, key="ri_employees")
                avg_salary_retention = st.number_input("Average Salary ($)", min_value=0, value=75000, step=5000, key="ri_salary")
                
                annual_savings_potential = (improvement_potential / 100) * total_employees * avg_salary_retention * 1.5
                st.success(f"**Annual Savings Potential: {format_currency(annual_savings_potential)}**")
        
        with tab4:
            st.subheader("💵 Revenue Impact Calculators")
            
            col1, col2 = st.columns(2)
            with col1:
                st.markdown("**Daily Revenue per Employee Calculator**")
                
                calculation_method = st.selectbox("Calculation Method", ["Revenue-based", "Profit-based", "Role-specific"], key="rev_method")
                
                if calculation_method == "Revenue-based":
                    annual_revenue = st.number_input("Annual Company Revenue ($)", min_value=0, value=10000000, step=100000, key="rev_annual")
                    total_employees_rev = st.number_input("Total Employees", min_value=1, value=100, key="rev_employees")
                    revenue_attribution = st.slider("Revenue Attribution (%)", 0, 100, 70, help="% of revenue attributable to employees vs. capital/systems", key="rev_attribution")
                    
                    daily_revenue_per_employee = (annual_revenue * (revenue_attribution/100)) / total_employees_rev / 250
                    st.success(f"**Daily Revenue per Employee: {format_currency(daily_revenue_per_employee)}**")
                
                elif calculation_method == "Role-specific":
                    role_type = st.selectbox("Role Type", ["Sales", "Customer Service", "Production", "Support"], key="rev_role")
                    if role_type == "Sales":
                        annual_quota = st.number_input("Annual Sales Quota ($)", min_value=0, value=500000, step=10000, key="rev_quota")
                        daily_revenue_per_employee = annual_quota / 250
                    else:
                        st.info("Select role-specific metrics based on your industry")
                        daily_revenue_per_employee = 800  # Default
                    
                    st.success(f"**Daily Revenue per Employee: {format_currency(daily_revenue_per_employee)}**")
            
            with col2:
                st.markdown("**Productivity Loss Calculator**")
                position_criticality = st.selectbox("Position Criticality", ["Non-critical", "Important", "Critical", "Mission-critical"], key="pl_criticality")
                criticality_loss = {"Non-critical": 30, "Important": 50, "Critical": 70, "Mission-critical": 90}
                
                baseline_loss = criticality_loss[position_criticality]
                
                backup_coverage = st.selectbox("Backup Coverage", ["None", "Partial", "Good", "Excellent"], key="pl_backup")
                coverage_reduction = {"None": 0, "Partial": 10, "Good": 20, "Excellent": 30}
                
                final_productivity_loss = max(10, baseline_loss - coverage_reduction[backup_coverage])
                
                st.success(f"**Estimated Productivity Loss: {final_productivity_loss}%**")
                
                st.markdown("**Customer Impact Calculator**")
                customer_facing = st.selectbox("Customer-Facing Role?", ["No", "Indirect", "Direct", "Primary contact"], key="ci_facing")
                impact_factors = {"No": 0, "Indirect": 5, "Direct": 15, "Primary contact": 25}
                
                customer_impact = impact_factors[customer_facing]
                st.success(f"**Customer Impact Factor: {customer_impact}%**")
        
        st.markdown("---")
        st.info("💡 **Tip:** Use these calculated values as inputs in your ROI analysis above. Save this page for future reference!")
    
    # Methodology explanation
    with st.expander("📊 About Our ROI Methodology", expanded=False):
        st.markdown("""
        **Incremental Cost Accounting Approach:**
        
        - **Costs**: Only true incremental expenses (facilitator fees, materials, venues, travel)
        - **Benefits**: Incremental value above baseline performance
        - **Excluded**: Employee salary costs (paid regardless of training participation)
        
        This approach provides a more accurate view of the actual investment required and expected returns,
        avoiding the "apples to oranges" comparison of treating salaries as costs while calculating productivity 
        gains as benefits.
        """)
    
    # Show export capabilities
    col1, col2, col3 = st.columns(3)
    with col1:
        if REPORTLAB_AVAILABLE:
            st.success("✅ PDF Export Available")
        else:
            st.warning("⚠️ PDF Export: Install reportlab")
    
    with col2:
        if PPTX_AVAILABLE:
            st.success("✅ PowerPoint Export Available")
        else:
            st.warning("⚠️ PowerPoint: Install python-pptx")
    
    with col3:
        st.success("✅ Text & JSON Export Available")
    
    # Installation instructions if needed
    if not REPORTLAB_AVAILABLE or not PPTX_AVAILABLE:
        with st.expander("📦 Enable Additional Export Options"):
            st.write("To enable all export formats, install the following packages:")
            if not REPORTLAB_AVAILABLE:
                st.code("pip install reportlab", language="bash")
                st.write("• **reportlab**: Enables PDF report generation with tables and charts")
            if not PPTX_AVAILABLE:
                st.code("pip install python-pptx", language="bash") 
                st.write("• **python-pptx**: Enables PowerPoint presentation generation")
            st.info("After installation, restart your Streamlit application to enable these features.")
    
    # Initialize session state
    if 'selected_initiatives' not in st.session_state:
        st.session_state.selected_initiatives = []
    if 'params' not in st.session_state:
        st.session_state.params = {}
    
    # Sidebar for initiative selection
    with st.sidebar:
        st.header("🎯 Select HR Initiatives")
        
        # Initiative Templates
        st.subheader("📋 Available Templates")
        
        for key, template in INITIATIVE_TEMPLATES.items():
            with st.expander(f"📊 {template['name']}"):
                st.write(f"**Description:** {template['description']}")
                st.write(f"**Typical ROI:** {template['typical_roi']}")
                
                if st.button(f"Add {template['name']}", key=f"add_{key}"):
                    if key not in st.session_state.selected_initiatives:
                        st.session_state.selected_initiatives.append(key)
                        st.session_state.params[key] = template.copy()
                        st.success(f"Added {template['name']}!")
                        st.rerun()
        
        st.divider()
        
        # Currently selected initiatives
        st.subheader("✅ Selected Initiatives")
        if st.session_state.selected_initiatives:
            st.info("💡 Click the red 🗑️ button next to any initiative to remove it")
            
            for initiative in st.session_state.selected_initiatives:
                col1, col2 = st.columns([4, 1])
                with col1:
                    st.write(f"📊 {INITIATIVE_TEMPLATES[initiative]['name']}")
                with col2:
                    if st.button("🗑️", key=f"remove_{initiative}", help="Remove this initiative", type="secondary"):
                        st.session_state.selected_initiatives.remove(initiative)
                        if initiative in st.session_state.params:
                            del st.session_state.params[initiative]
                        st.success(f"✅ Removed: {INITIATIVE_TEMPLATES[initiative]['name']}")
                        st.rerun()
            
            # Clear all button
            st.divider()
            col1, col2 = st.columns(2)
            with col1:
                if st.button("🗑️ Remove All", help="Remove all selected initiatives", type="secondary"):
                    cleared_count = len(st.session_state.selected_initiatives)
                    st.session_state.selected_initiatives.clear()
                    st.session_state.params.clear()
                    st.success(f"✅ Cleared {cleared_count} initiatives!")
                    st.rerun()
            with col2:
                st.write(f"**Total selected:** {len(st.session_state.selected_initiatives)}")
        else:
            st.info("👈 No initiatives selected yet. Use the 'Add' buttons above to get started!")
            st.markdown("**How to remove initiatives:**")
            st.markdown("- Individual: Click 🗑️ next to each initiative")
            st.markdown("- All at once: Use 'Remove All' button")
    
    # Main content
    if not st.session_state.selected_initiatives:
        st.info("👈 Please select one or more HR initiatives from the sidebar to begin calculating ROI.")
        return
    
    # Create tabs for each selected initiative
    if len(st.session_state.selected_initiatives) == 1:
        # Single initiative - no tabs needed
        initiative_key = st.session_state.selected_initiatives[0]
        display_initiative(initiative_key)
    else:
        # Multiple initiatives - create tabs
        tab_names = [INITIATIVE_TEMPLATES[key]['name'] for key in st.session_state.selected_initiatives]
        tabs = st.tabs(tab_names)
        
        for i, initiative_key in enumerate(st.session_state.selected_initiatives):
            with tabs[i]:
                display_initiative(initiative_key)
    
    # Overall summary if multiple initiatives
    if len(st.session_state.selected_initiatives) > 1:
        st.divider()
        display_overall_summary()

def display_initiative(initiative_key):
    """Display interface for a specific initiative"""
    template = INITIATIVE_TEMPLATES[initiative_key]
    params = st.session_state.params[initiative_key]
    
    st.subheader(f"📊 {template['name']}")
    st.write(template['description'])
    
    # Parameters input section
    with st.expander("⚙️ Adjust Parameters", expanded=True):
        col1, col2 = st.columns(2)
        
        # Update parameters based on initiative type
        if initiative_key == 'leadership_development' or initiative_key == 'executive_coaching':
            with col1:
                st.markdown("**📊 Program Parameters**")
                params['participants'] = st.number_input(
                    "Number of Participants", 
                    min_value=1, 
                    value=params['participants'],
                    key=f"participants_{initiative_key}"
                )
                params['avg_salary'] = st.number_input(
                    "Average Salary ($)", 
                    min_value=0, 
                    value=params['avg_salary'], 
                    step=5000,
                    key=f"salary_{initiative_key}"
                )
                params['program_duration'] = st.number_input(
                    "Program Duration (months)", 
                    min_value=1, 
                    value=params['program_duration'],
                    key=f"duration_{initiative_key}"
                )
                
                st.markdown("**💰 Direct Costs**")
                params['facilitator_costs'] = st.number_input(
                    "Facilitator Costs ($)", 
                    min_value=0, 
                    value=params['facilitator_costs'], 
                    step=5000,
                    key=f"facilitator_{initiative_key}"
                )
                params['materials_costs'] = st.number_input(
                    "Materials Costs ($)", 
                    min_value=0, 
                    value=params['materials_costs'], 
                    step=1000,
                    key=f"materials_{initiative_key}"
                )
                params['venue_costs'] = st.number_input(
                    "Venue Costs ($)", 
                    min_value=0, 
                    value=params['venue_costs'], 
                    step=1000,
                    key=f"venue_{initiative_key}"
                )
                params['travel_costs'] = st.number_input(
                    "Travel Costs ($)", 
                    min_value=0, 
                    value=params.get('travel_costs', 20000), 
                    step=1000,
                    key=f"travel_{initiative_key}"
                )
            
            with col2:
                st.markdown("**📈 Expected Improvements**")
                params['productivity_gain'] = st.slider(
                    "Productivity Improvement (%)", 
                    0, 50, 
                    params['productivity_gain'],
                    help="Expected increase in individual productivity",
                    key=f"productivity_{initiative_key}"
                )
                params['retention_improvement'] = st.slider(
                    "Retention Improvement (%)", 
                    0, 50, 
                    params['retention_improvement'],
                    help="Reduction in turnover rate for participants",
                    key=f"retention_{initiative_key}"
                )
                params['team_performance_gain'] = st.slider(
                    "Team Performance Gain (%)", 
                    0, 30, 
                    params['team_performance_gain'],
                    help="Improvement in team performance led by participants",
                    key=f"team_{initiative_key}"
                )
                params['sick_leave_reduction'] = st.slider(
                    "Sick Leave Reduction (%)", 
                    0, 40, 
                    params.get('sick_leave_reduction', 20),
                    help="Reduction in sick days due to better leadership and work environment",
                    key=f"sick_leave_{initiative_key}"
                )
                
                st.markdown("**⚙️ Advanced Settings**")
                params['current_turnover'] = st.number_input(
                    "Current Turnover Rate (%)", 
                    min_value=0.0, 
                    max_value=50.0,
                    value=params.get('current_turnover', 18.0), 
                    step=1.0,
                    key=f"turnover_{initiative_key}"
                )
                params['team_size'] = st.number_input(
                    "Average Team Size", 
                    min_value=1, 
                    value=params.get('team_size', 8),
                    key=f"teamsize_{initiative_key}"
                )
                params['current_sick_days'] = st.number_input(
                    "Current Sick Days per Employee/Year", 
                    min_value=0, 
                    value=params.get('current_sick_days', 8),
                    help="Average sick days taken per employee annually",
                    key=f"sickdays_{initiative_key}"
                )
            
            # Calculate and display results
            results = calculate_leadership_roi(params)
            
        elif initiative_key == 'time_to_fill_optimization':
            with col1:
                st.markdown("**📊 Current State**")
                params['annual_positions'] = st.number_input(
                    "Annual Positions to Fill", 
                    min_value=1, 
                    value=params['annual_positions'],
                    help="Number of positions that need to be filled annually",
                    key=f"positions_{initiative_key}"
                )
                params['current_time_to_fill'] = st.number_input(
                    "Current Time to Fill (days)", 
                    min_value=1, 
                    value=params['current_time_to_fill'],
                    help="Average days from job posting to offer acceptance",
                    key=f"current_time_{initiative_key}"
                )
                params['target_time_to_fill'] = st.number_input(
                    "Target Time to Fill (days)", 
                    min_value=1, 
                    value=params['target_time_to_fill'],
                    help="Goal for average days to fill positions",
                    key=f"target_time_{initiative_key}"
                )
                params['avg_position_salary'] = st.number_input(
                    "Average Position Salary ($)", 
                    min_value=0, 
                    value=params['avg_position_salary'], 
                    step=5000,
                    key=f"avg_salary_{initiative_key}"
                )
                
                st.markdown("**💰 Investment Costs**")
                params['optimization_investment'] = st.number_input(
                    "Process Optimization Investment ($)", 
                    min_value=0, 
                    value=params['optimization_investment'], 
                    step=5000,
                    help="Investment in process improvements and consulting",
                    key=f"optimization_{initiative_key}"
                )
                params['training_costs'] = st.number_input(
                    "Recruiter Training Costs ($)", 
                    min_value=0, 
                    value=params.get('training_costs', 15000), 
                    step=1000,
                    key=f"training_{initiative_key}"
                )
                params['technology_costs'] = st.number_input(
                    "Technology & Tools ($)", 
                    min_value=0, 
                    value=params.get('technology_costs', 25000), 
                    step=1000,
                    help="ATS upgrades, automation tools, etc.",
                    key=f"technology_{initiative_key}"
                )
            
            with col2:
                st.markdown("**📈 Cost Impact Parameters**")
                params['productivity_loss_rate'] = st.slider(
                    "Productivity Loss During Vacancy (%)", 
                    0, 100, 
                    int(params['productivity_loss_rate']),
                    help="% of position's productivity lost while vacant",
                    key=f"prod_loss_{initiative_key}"
                )
                params['team_impact_factor'] = st.slider(
                    "Team Productivity Impact (%)", 
                    0, 30, 
                    int(params['team_impact_factor']),
                    help="% productivity hit on team due to vacancy stress",
                    key=f"team_impact_{initiative_key}"
                )
                params['overtime_multiplier'] = st.number_input(
                    "Overtime Rate Multiplier", 
                    min_value=1.0, 
                    max_value=3.0,
                    value=params['overtime_multiplier'], 
                    step=0.1,
                    help="Overtime pay rate (1.5 = time and a half)",
                    key=f"overtime_{initiative_key}"
                )
                
                st.markdown("**💵 Revenue Impact Parameters**")
                params['revenue_generating_percentage'] = st.slider(
                    "Revenue-Generating Roles (%)", 
                    0, 100, 
                    int(params.get('revenue_generating_percentage', 60)),
                    help="% of positions that directly generate revenue",
                    key=f"revenue_pct_{initiative_key}"
                )
                params['revenue_per_employee_daily'] = st.number_input(
                    "Daily Revenue per Employee ($)", 
                    min_value=0, 
                    value=params.get('revenue_per_employee_daily', 800),
                    step=50,
                    help="Average daily revenue generated per employee",
                    key=f"daily_revenue_{initiative_key}"
                )
                params['customer_impact_factor'] = st.slider(
                    "Customer Service Impact (%)", 
                    0, 50, 
                    int(params.get('customer_impact_factor', 25)),
                    help="% revenue impact due to reduced service quality during vacancies",
                    key=f"customer_impact_{initiative_key}"
                )
                
                st.markdown("**⚙️ Advanced Settings**")
                params['team_size'] = st.number_input(
                    "Average Team Size", 
                    min_value=1, 
                    value=params.get('team_size', 6),
                    help="Number of team members affected by each vacancy",
                    key=f"teamsize_ttf_{initiative_key}"
                )
                params['overtime_hours_per_day'] = st.number_input(
                    "Overtime Hours per Day", 
                    min_value=0, 
                    value=params.get('overtime_hours_per_day', 2),
                    help="Daily overtime hours needed to cover vacant position",
                    key=f"ot_hours_{initiative_key}"
                )
            
            results = calculate_time_to_fill_roi(params)
            
        elif initiative_key == 'onboarding_excellence':
            with col1:
                params['annual_new_hires'] = st.number_input(
                    "Annual New Hires", 
                    min_value=1, 
                    value=params['annual_new_hires'],
                    key=f"new_hires_{initiative_key}"
                )
                params['current_time_to_productivity'] = st.number_input(
                    "Current Time to Productivity (months)", 
                    min_value=1, 
                    value=params['current_time_to_productivity'],
                    key=f"productivity_time_{initiative_key}"
                )
            
            with col2:
                params['productivity_acceleration'] = st.slider(
                    "Productivity Acceleration (%)", 
                    0, 70, 
                    params['productivity_acceleration'],
                    key=f"acceleration_{initiative_key}"
                )
                params['onboarding_retention_improvement'] = st.slider(
                    "Retention Improvement (%)", 
                    0, 40, 
                    params['onboarding_retention_improvement'],
                    key=f"onboard_retention_{initiative_key}"
                )
            
            results = calculate_onboarding_roi(params)
            
        elif initiative_key == 'engagement_retention':
            with col1:
                params['total_employees'] = st.number_input(
                    "Total Employees", 
                    min_value=1, 
                    value=params['total_employees'],
                    key=f"employees_{initiative_key}"
                )
                params['current_engagement_score'] = st.number_input(
                    "Current Engagement Score (1-10)", 
                    min_value=1.0, 
                    max_value=10.0,
                    value=params['current_engagement_score'],
                    step=0.1,
                    key=f"engagement_{initiative_key}"
                )
            
            with col2:
                params['engagement_improvement'] = st.slider(
                    "Engagement Score Improvement", 
                    0.0, 3.0, 
                    params['engagement_improvement'],
                    step=0.1,
                    key=f"engagement_improve_{initiative_key}"
                )
                params['retention_improvement'] = st.slider(
                    "Retention Improvement (%)", 
                    0, 50, 
                    params['retention_improvement'],
                    key=f"retention_improve_{initiative_key}"
                )
            
            results = calculate_engagement_roi(params)
            
        elif initiative_key == 'talent_development':
            with col1:
                params['development_participants'] = st.number_input(
                    "Development Participants", 
                    min_value=1, 
                    value=params['development_participants'],
                    key=f"dev_participants_{initiative_key}"
                )
                params['internal_mobility_rate'] = st.number_input(
                    "Current Internal Mobility Rate (%)", 
                    min_value=0, 
                    max_value=100,
                    value=params['internal_mobility_rate'],
                    key=f"mobility_{initiative_key}"
                )
            
            with col2:
                params['mobility_improvement'] = st.slider(
                    "Mobility Improvement (%)", 
                    0, 50, 
                    params['mobility_improvement'],
                    key=f"mobility_improve_{initiative_key}"
                )
                params['succession_improvement'] = st.slider(
                    "Succession Planning Improvement (%)", 
                    0, 50, 
                    params['succession_improvement'],
                    key=f"succession_{initiative_key}"
                )
            
            results = calculate_development_roi(params)
    
    # Display results
    st.subheader("📈 Results")
    
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        st.metric(
            "ROI",
            f"{results['roi']:.0f}%",
            delta=get_roi_status(results['roi'])
        )
    
    with col2:
        investment_key = 'total_investment' if 'total_investment' in results else 'total_costs'
        st.metric(
            "Incremental Investment",
            format_currency(results[investment_key])
        )
    
    with col3:
        benefits_key = 'annual_benefits' if 'annual_benefits' in results else 'annual_savings'
        st.metric(
            "Annual Benefits",
            format_currency(results[benefits_key])
        )
    
    with col4:
        if 'payback_months' in results:
            st.metric(
                "Payback Period",
                f"{results['payback_months']:.1f} months"
            )
        else:
            net_benefit = results[benefits_key] - results[investment_key]
            st.metric(
                "Net Annual Benefit",
                format_currency(net_benefit)
            )
    
    # Special metrics for time to fill
    if initiative_key == 'time_to_fill_optimization':
        # Time improvement metrics
        st.subheader("⏱️ Time to Fill Improvements")
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            st.metric(
                "Days Saved per Position",
                f"{results['days_saved_per_position']:.0f} days"
            )
        
        with col2:
            st.metric(
                "Total Days Saved Annually",
                f"{results['total_days_saved_annually']:.0f} days"
            )
        
        with col3:
            current_time = params.get('current_time_to_hire', 45)
            target_time = params.get('target_time_to_fill', 35)
            improvement_pct = ((current_time - target_time) / current_time) * 100
            st.metric(
                "Time Reduction",
                f"{improvement_pct:.0f}%",
                delta=f"{current_time}→{target_time} days"
            )
        
        with col4:
            if 'payback_months' in results:
                st.metric(
                    "Payback Period",
                    f"{results['payback_months']:.1f} months"
                )
        
        # Revenue at risk analysis
        st.subheader("💰 Revenue at Risk Analysis")
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            daily_revenue_at_risk = results.get('daily_revenue_at_risk', 0)
            st.metric(
                "Daily Revenue at Risk",
                format_currency(daily_revenue_at_risk),
                help="Revenue at risk per day with current time to fill"
            )
        
        with col2:
            st.metric(
                "Revenue-Generating Positions",
                f"{results.get('revenue_positions', 0):.0f}",
                delta=f"{params.get('revenue_generating_percentage', 60)}% of total"
            )
        
        with col3:
            total_revenue_protection = results.get('total_revenue_protection', 0)
            st.metric(
                "Annual Revenue Protection",
                format_currency(total_revenue_protection),
                help="Revenue protected by faster time to fill"
            )
        
        with col4:
            current_revenue_risk = results.get('total_revenue_at_risk_current', 1)
            target_revenue_risk = results.get('total_revenue_at_risk_target', 0)
            revenue_risk_reduction = ((current_revenue_risk - target_revenue_risk) / current_revenue_risk) * 100 if current_revenue_risk > 0 else 0
            st.metric(
                "Revenue Risk Reduction",
                f"{revenue_risk_reduction:.0f}%",
                help="Reduction in revenue at risk per position"
            )
        
        # Dual analysis breakdown
        st.subheader("📊 Cost Savings vs Revenue Protection Analysis")
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("**💸 Cost Savings Breakdown**")
            cost_breakdown = results.get('cost_savings_breakdown', {})
            cost_breakdown_df = pd.DataFrame([
                {"Category": "Productivity Recovery", "Annual Value": cost_breakdown.get('productivity_recovery', 0)},
                {"Category": "Overtime Reduction", "Annual Value": cost_breakdown.get('overtime_savings', 0)},
                {"Category": "Team Productivity", "Annual Value": cost_breakdown.get('team_productivity_gain', 0)},
                {"Category": "Faster Onboarding", "Annual Value": cost_breakdown.get('faster_onboarding', 0)},
            ])
            cost_breakdown_df['Annual Value'] = cost_breakdown_df['Annual Value'].apply(format_currency)
            st.dataframe(cost_breakdown_df, hide_index=True, use_container_width=True)
            
            st.metric("Total Cost Savings", format_currency(results.get('total_cost_savings', 0)))
        
        with col2:
            st.markdown("**💵 Revenue Protection Breakdown**")
            revenue_breakdown = results.get('revenue_impact_breakdown', {})
            revenue_breakdown_df = pd.DataFrame([
                {"Category": "Direct Revenue Protection", "Annual Value": revenue_breakdown.get('direct_revenue_protection', 0)},
                {"Category": "Customer Service Impact", "Annual Value": revenue_breakdown.get('customer_impact_protection', 0)},
                {"Category": "Opportunity Cost Protection", "Annual Value": revenue_breakdown.get('opportunity_cost_protection', 0)},
                {"Category": "Market Share Protection", "Annual Value": revenue_breakdown.get('market_share_protection', 0)},
            ])
            revenue_breakdown_df['Annual Value'] = revenue_breakdown_df['Annual Value'].apply(format_currency)
            st.dataframe(revenue_breakdown_df, hide_index=True, use_container_width=True)
            
            st.metric("Total Revenue Protection", format_currency(results.get('total_revenue_protection', 0)))
        
        # Combined chart
        st.subheader("📈 Combined Benefits Analysis")
        
        # Create combined breakdown for chart with safe access
        cost_savings = results.get('cost_savings_breakdown', {})
        revenue_impact = results.get('revenue_impact_breakdown', {})
        
        combined_benefits = {}
        for k, v in cost_savings.items():
            combined_benefits[f"Cost: {k}"] = v
        for k, v in revenue_impact.items():
            combined_benefits[f"Revenue: {k}"] = v
        
        if combined_benefits:  # Only create chart if we have data
            fig = px.bar(
                x=list(combined_benefits.keys()),
                y=list(combined_benefits.values()),
                title="Time to Fill Optimization - Cost Savings vs Revenue Protection",
                color=['Cost Savings'] * len(cost_savings) + ['Revenue Protection'] * len(revenue_impact),
                color_discrete_map={'Cost Savings': '#1f77b4', 'Revenue Protection': '#2ca02c'}
            )
            fig.update_layout(
                xaxis_title="Benefit Category",
                yaxis_title="Annual Value ($)",
                xaxis_tickangle=-45,
                showlegend=True
            )
            st.plotly_chart(fig, use_container_width=True)
    
    # Cost breakdown for leadership/coaching programs
    if initiative_key in ['leadership_development', 'executive_coaching'] and 'cost_breakdown' in results:
        col1, col2 = st.columns(2)
        
        with col1:
            st.subheader("💰 Cost Breakdown")
            cost_breakdown = results['cost_breakdown']
            cost_df = pd.DataFrame([
                {"Cost Category": "Facilitator", "Amount": cost_breakdown['facilitator_costs']},
                {"Cost Category": "Materials", "Amount": cost_breakdown['materials_costs']},
                {"Cost Category": "Venue", "Amount": cost_breakdown['venue_costs']},
                {"Cost Category": "Travel", "Amount": cost_breakdown['travel_costs']},
            ])
            st.dataframe(cost_df, hide_index=True)
            
        with col2:
            # Cost pie chart
            fig_cost = px.pie(
                values=list(cost_breakdown.values()),
                names=list(cost_breakdown.keys()),
                title="Investment Breakdown"
            )
            st.plotly_chart(fig_cost, use_container_width=True)
    
    # Investment breakdown for time to fill optimization
    if initiative_key == 'time_to_fill_optimization' and 'investment_breakdown' in results:
        col1, col2 = st.columns(2)
        
        with col1:
            st.subheader("💰 Investment Breakdown")
            investment_breakdown = results['investment_breakdown']
            investment_df = pd.DataFrame([
                {"Investment Category": "Process Optimization", "Amount": investment_breakdown['optimization_investment']},
                {"Investment Category": "Training", "Amount": investment_breakdown['training_costs']},
                {"Investment Category": "Technology", "Amount": investment_breakdown['technology_costs']},
            ])
            st.dataframe(investment_df, hide_index=True)
            
        with col2:
            # Investment pie chart
            fig_investment = px.pie(
                values=list(investment_breakdown.values()),
                names=["Process Optimization", "Training", "Technology"],
                title="Investment Breakdown"
            )
            st.plotly_chart(fig_investment, use_container_width=True)
    
    # Benefits breakdown chart
    if 'benefit_breakdown' in results:
        breakdown = results['benefit_breakdown']
        
        fig = px.bar(
            x=list(breakdown.keys()),
            y=list(breakdown.values()),
            title=f"{template['name']} - Annual Benefits Breakdown",
            color=list(breakdown.values()),
            color_continuous_scale="Viridis"
        )
        fig.update_layout(
            xaxis_title="Benefit Category",
            yaxis_title="Annual Value ($)",
            showlegend=False
        )
        st.plotly_chart(fig, use_container_width=True)
    elif 'cost_savings_breakdown' in results and 'revenue_impact_breakdown' in results:
        # This is handled in the time to fill section above
        pass
    
    # Export options for individual initiative
    st.divider()
    st.subheader("📄 Export Options")
    
    col1, col2, col3 = st.columns(3)
    
    with col1:
        if st.button(f"📋 Text Summary", key=f"export_text_{initiative_key}"):
            individual_report = f"""
{template['name']} - ROI Analysis
Generated: {datetime.now().strftime('%B %d, %Y')}
Methodology: Incremental Cost Accounting

RESULTS SUMMARY
===============
ROI: {results['roi']:.0f}%
Incremental Investment: {format_currency(results.get('total_investment', results.get('total_costs', 0)))}
Annual Benefits: {format_currency(results.get('annual_benefits', results.get('annual_savings', 0)))}
Net Annual Benefit: {format_currency(results.get('net_annual_benefit', results.get('annual_benefits', results.get('annual_savings', 0)) - results.get('total_investment', results.get('total_costs', 0))))}
Status: {get_roi_status(results['roi'])}

PARAMETERS USED
===============
"""
            for key, value in params.items():
                if isinstance(value, (int, float)):
                    individual_report += f"{key.replace('_', ' ').title()}: {value:,}\n"
                else:
                    individual_report += f"{key.replace('_', ' ').title()}: {value}\n"
            
            st.download_button(
                label="📥 Download Text",
                data=individual_report,
                file_name=f"{initiative_key}_roi_summary_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                mime="text/plain",
                key=f"download_{initiative_key}"
            )
    
    with col2:
        if REPORTLAB_AVAILABLE:
            if st.button(f"📄 PDF Report", key=f"export_pdf_{initiative_key}"):
                # Create single initiative data structure for PDF
                investment = results.get('total_investment', results.get('total_costs', 0))
                benefits = results.get('annual_benefits', results.get('annual_savings', 0))
                
                single_initiative_data = [{
                    'Initiative': template['name'],
                    'Investment': investment,
                    'Annual Benefits': benefits,
                    'ROI (%)': results['roi']
                }]
                
                pdf_buffer = create_pdf_report(
                    single_initiative_data, 
                    results['roi'], 
                    investment, 
                    benefits, 
                    {initiative_key: params}
                )
                if pdf_buffer:
                    st.download_button(
                        label="📥 Download PDF",
                        data=pdf_buffer,
                        file_name=f"{initiative_key}_roi_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf",
                        mime="application/pdf",
                        key=f"download_pdf_{initiative_key}"
                    )
        else:
            st.button(f"📄 PDF Report", disabled=True, help="Install reportlab to enable PDF export")
    
    with col3:
        if st.button(f"📊 JSON Data", key=f"export_json_{initiative_key}"):
            investment = results.get('total_investment', results.get('total_costs', 0))
            benefits = results.get('annual_benefits', results.get('annual_savings', 0))
            
            export_data = {
                'methodology': 'incremental_cost_accounting',
                'initiative': {
                    'name': template['name'],
                    'investment': investment,
                    'annual_benefits': benefits,
                    'roi': results['roi'],
                    'results': results
                },
                'parameters': params,
                'timestamp': datetime.now().isoformat()
            }
            st.download_button(
                label="📥 Download JSON",
                data=json.dumps(export_data, indent=2, default=str),
                file_name=f"{initiative_key}_roi_data_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json",
                mime="application/json",
                key=f"download_json_{initiative_key}"
            )
    
    col1, col2 = st.columns(2)
    with col1:
        st.info(f"💡 **Expected ROI Range:** {template['typical_roi']}")
        
    with col2:
        # Special notes for different calculators
        if initiative_key == 'time_to_fill_optimization':
            st.info("""
            **💡 Business Impact Analysis:**
            
            **Cost Savings:**
            • Productivity recovery (faster placement)
            • Overtime cost reduction  
            • Team productivity improvement
            • Faster new hire value realization
            
            **Revenue Protection:**
            • Direct revenue from vacant positions
            • Customer service impact mitigation
            • Opportunity cost protection
            • Market share protection
            """)
        else:
            st.info(f"💡 **Expected ROI Range:** {template['typical_roi']}")

def display_overall_summary():
    """Display summary across all selected initiatives"""
    st.subheader("🎯 Overall Portfolio Summary")
    
    total_investment = 0
    total_benefits = 0
    initiative_results = []
    
    for initiative_key in st.session_state.selected_initiatives:
        params = st.session_state.params[initiative_key]
        
        if initiative_key in ['leadership_development', 'executive_coaching']:
            results = calculate_leadership_roi(params)
            investment = results['total_costs']
            benefits = results['annual_benefits']
        elif initiative_key == 'time_to_fill_optimization':
            results = calculate_time_to_fill_roi(params)
            investment = results['total_investment']
            benefits = results['annual_benefits']
        elif initiative_key == 'onboarding_excellence':
            results = calculate_onboarding_roi(params)
            investment = results['total_investment']
            benefits = results['annual_benefits']
        elif initiative_key == 'engagement_retention':
            results = calculate_engagement_roi(params)
            investment = results['total_investment']
            benefits = results['annual_benefits']
        elif initiative_key == 'talent_development':
            results = calculate_development_roi(params)
            investment = results['total_investment']
            benefits = results['annual_benefits']
        
        total_investment += investment
        total_benefits += benefits
        
        initiative_results.append({
            'Initiative': INITIATIVE_TEMPLATES[initiative_key]['name'],
            'Investment': investment,
            'Annual Benefits': benefits,
            'ROI (%)': results['roi']
        })
    
    overall_roi = ((total_benefits - total_investment) / total_investment * 100) if total_investment > 0 else 0
    
    # Overall metrics
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        st.metric("Portfolio ROI", f"{overall_roi:.0f}%", delta=get_roi_status(overall_roi))
    
    with col2:
        st.metric("Total Incremental Investment", format_currency(total_investment))
    
    with col3:
        st.metric("Total Annual Benefits", format_currency(total_benefits))
    
    with col4:
        st.metric("Net Annual Benefit", format_currency(total_benefits - total_investment))
    
    # Initiative comparison
    df = pd.DataFrame(initiative_results)
    df = df.sort_values('ROI (%)', ascending=False)
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.subheader("📊 Initiative Comparison")
        # Format the dataframe for better display
        display_df = df.copy()
        display_df['Investment'] = display_df['Investment'].apply(format_currency)
        display_df['Annual Benefits'] = display_df['Annual Benefits'].apply(format_currency)
        display_df['ROI (%)'] = display_df['ROI (%)'].apply(lambda x: f"{x:.0f}%")
        st.dataframe(display_df, use_container_width=True, hide_index=True)
    
    with col2:
        fig = px.bar(
            df,
            x='Initiative',
            y='ROI (%)',
            title="ROI by Initiative",
            color='ROI (%)',
            color_continuous_scale="RdYlGn"
        )
        fig.update_layout(xaxis_tickangle=-45)
        st.plotly_chart(fig, use_container_width=True)
    
    # Export functionality
    st.divider()
    st.subheader("📄 Export Options")
    
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        if st.button("📋 Text Report", type="primary"):
            report = create_summary_report(initiative_results, overall_roi, total_investment, total_benefits)
            st.download_button(
                label="📥 Download Text Report",
                data=report,
                file_name=f"hr_roi_summary_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                mime="text/plain"
            )
    
    with col2:
        if REPORTLAB_AVAILABLE:
            if st.button("📄 PDF Report", type="primary"):
                pdf_buffer = create_pdf_report(initiative_results, overall_roi, total_investment, total_benefits, st.session_state.params)
                if pdf_buffer:
                    st.download_button(
                        label="📥 Download PDF",
                        data=pdf_buffer,
                        file_name=f"hr_roi_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf",
                        mime="application/pdf"
                    )
        else:
            st.button("📄 PDF Report", disabled=True, help="Install reportlab to enable PDF export")
    
    with col3:
        if PPTX_AVAILABLE:
            if st.button("📊 PowerPoint", type="primary"):
                ppt_buffer = create_powerpoint_presentation(initiative_results, overall_roi, total_investment, total_benefits)
                if ppt_buffer:
                    st.download_button(
                        label="📥 Download PowerPoint",
                        data=ppt_buffer,
                        file_name=f"hr_roi_presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
        else:
            st.button("📊 PowerPoint", disabled=True, help="Install python-pptx to enable PowerPoint export")
    
    with col4:
        if st.button("📊 JSON Data", type="secondary"):
            export_data = {
                'methodology': 'incremental_cost_accounting',
                'initiatives': initiative_results,
                'summary': {
                    'total_investment': total_investment,
                    'total_benefits': total_benefits,
                    'overall_roi': overall_roi
                },
                'parameters': st.session_state.params,
                'timestamp': datetime.now().isoformat()
            }
            st.download_button(
                label="📥 Download JSON",
                data=json.dumps(export_data, indent=2, default=str),
                file_name=f"hr_roi_data_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json",
                mime="application/json"
            )

def create_summary_report(initiative_results, overall_roi, total_investment, total_benefits):
    """Create a summary report"""
    report = f"""
HR ROI CALCULATOR - SUMMARY REPORT (INCREMENTAL COST METHOD)
Generated: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}

METHODOLOGY
===========
This analysis uses incremental cost accounting:
- Costs: Only true additional expenses (facilitator, materials, venues, travel)
- Benefits: Incremental value above baseline performance  
- Excluded: Employee salary costs (paid regardless of training participation)

EXECUTIVE SUMMARY
================
Portfolio ROI: {overall_roi:.0f}%
Total Incremental Investment: {format_currency(total_investment)}
Total Annual Benefits: {format_currency(total_benefits)}
Net Annual Benefit: {format_currency(total_benefits - total_investment)}

Status: {get_roi_status(overall_roi)}

INITIATIVE BREAKDOWN
===================
"""
    
    for initiative in initiative_results:
        report += f"""
{initiative['Initiative']}:
  Incremental Investment: {format_currency(initiative['Investment'])}
  Annual Benefits: {format_currency(initiative['Annual Benefits'])}
  ROI: {initiative['ROI (%)']:.0f}%
"""
    
    report += f"""

RECOMMENDATIONS
===============
{"✅ Exceptional portfolio - proceed with full implementation" if overall_roi >= 500 else "✅ Excellent portfolio - proceed with implementation, prioritize by ROI" if overall_roi >= 300 else "✅ Strong portfolio - proceed with implementation, prioritize by ROI" if overall_roi >= 200 else "⚠️ Review highest-performing initiatives for priority implementation" if overall_roi >= 100 else "❌ Reassess assumptions and focus on highest ROI initiatives only"}

Generated by HR ROI Calculator (Incremental Cost Method)
"""
    
    return report

if __name__ == "__main__":
    main()
